from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Brand Colors ────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0A, 0x0C, 0x0F)   # #0a0c0f
PANEL     = RGBColor(0x14, 0x18, 0x20)   # #141820
BORDER    = RGBColor(0x25, 0x2B, 0x35)   # #252b35
AMBER     = RGBColor(0xE8, 0xA0, 0x34)   # #e8a034
YELLOW    = RGBColor(0xFF, 0xE6, 0x00)   # #ffe600
WHITE     = RGBColor(0xF5, 0xF7, 0xFB)   # #f5f7fb
GRAY      = RGBColor(0x9B, 0xA5, 0xB3)   # #9ba5b3
CYAN      = RGBColor(0x22, 0xD3, 0xEE)   # #22d3ee
GREEN     = RGBColor(0x10, 0xB9, 0x81)   # #10b981
RED       = RGBColor(0xEF, 0x44, 0x44)   # #ef4444
BLUE      = RGBColor(0x3B, 0x82, 0xF6)   # #3b82f6
DARK_GRAY = RGBColor(0x44, 0x44, 0x55)   # #444455

W = Inches(13.33)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill=PANEL, line=None, line_w=Pt(1)):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_txt(slide, text, l, t, w, h, size=Pt(14), bold=False, color=WHITE,
            align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_para(tf, text, size=Pt(13), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, space_before=Pt(4), italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def add_label(slide, text, l, t, color=AMBER):
    add_txt(slide, text, l, t, Inches(4), Inches(0.3),
            size=Pt(8), bold=True, color=color)


def amber_bar(slide, w=W):
    add_rect(slide, 0, 0, w, Inches(0.07), fill=AMBER)


def section_chip(slide, label, l, t):
    add_rect(slide, l, t, Inches(1.8), Inches(0.28), fill=RGBColor(0x1C, 0x20, 0x2C), line=AMBER)
    add_txt(slide, label, l + Inches(0.1), t + Inches(0.02),
            Inches(1.6), Inches(0.25), size=Pt(7.5), bold=True, color=AMBER)


def bullet_block(slide, title, bullets, l, t, w, h,
                 title_color=CYAN, bullet_color=GRAY, title_size=Pt(11)):
    add_rect(slide, l, t, w, h, fill=PANEL, line=BORDER)
    add_txt(slide, title, l + Inches(0.15), t + Inches(0.12),
            w - Inches(0.3), Inches(0.28), size=title_size, bold=True, color=title_color)
    txb = slide.shapes.add_textbox(l + Inches(0.15), t + Inches(0.42),
                                   w - Inches(0.3), h - Inches(0.55))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = b
        run.font.size = Pt(10.5)
        run.font.color.rgb = bullet_color


prs = new_prs()

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
amber_bar(sl)

# Left accent stripe
add_rect(sl, 0, Inches(0.07), Inches(0.06), H - Inches(0.07), fill=AMBER)

# Big name
add_txt(sl, "DRISHTI",
        Inches(0.5), Inches(1.6), Inches(7), Inches(1.5),
        size=Pt(82), bold=True, color=WHITE)

# Amber underline
add_rect(sl, Inches(0.5), Inches(3.05), Inches(4.5), Inches(0.06), fill=AMBER)

add_txt(sl, "Dynamic Resource Intelligence for\nSmart Highway and Traffic Intervention",
        Inches(0.5), Inches(3.2), Inches(8), Inches(0.9),
        size=Pt(17), bold=False, color=GRAY)

add_txt(sl, "AI-Powered Traffic Operations Platform  ·  Bengaluru Traffic Police",
        Inches(0.5), Inches(4.2), Inches(9), Inches(0.4),
        size=Pt(12), bold=False, color=AMBER)

# Right side visual block
add_rect(sl, Inches(9.2), Inches(0.5), Inches(3.6), Inches(6.5), fill=PANEL, line=BORDER)
add_txt(sl, "LIVE SYSTEM", Inches(9.5), Inches(0.8), Inches(3), Inches(0.3),
        size=Pt(8), bold=True, color=AMBER)

stats = [
    ("10.2M", "Vehicles in Bengaluru"),
    ("1,500+", "New vehicles daily"),
    ("3 Layers", "Citizen · Ops · Field"),
    ("397", "ML Feature Dimensions"),
    ("2-Gate", "AI Firewall"),
    ("30s", "GPS Refresh Cycle"),
]
for i, (num, label) in enumerate(stats):
    y = Inches(1.3) + i * Inches(0.9)
    add_txt(sl, num, Inches(9.4), y, Inches(1.5), Inches(0.4),
            size=Pt(22), bold=True, color=CYAN)
    add_txt(sl, label, Inches(10.9), y + Inches(0.06), Inches(1.8), Inches(0.35),
            size=Pt(10), color=GRAY)

add_txt(sl, "Gridathon 2026  ·  Team DRISHTI",
        Inches(0.5), Inches(6.9), Inches(6), Inches(0.35),
        size=Pt(9), color=DARK_GRAY)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ════════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
amber_bar(sl)

section_chip(sl, "THE PROBLEM", Inches(0.4), Inches(0.2))
add_txt(sl, "Bengaluru's Traffic Crisis Has No Digital Backbone",
        Inches(0.4), Inches(0.6), Inches(10), Inches(0.7),
        size=Pt(28), bold=True, color=WHITE)
add_rect(sl, Inches(0.4), Inches(1.25), Inches(3.5), Inches(0.04), fill=AMBER)

# 3 problem cards
problems = [
    (RED,  "NO UNIFIED INTAKE",
     "Incident reports arrive via phone calls, WhatsApp messages, walk-ins, and radio — with no central log, no timestamp, and no priority system. The control room runs on memory and instinct."),
    (AMBER,"NO REAL-TIME VISIBILITY",
     "Officers in the field have no GPS tracking. The control room cannot see where units are deployed, which complaints are active, or which junctions are currently unmanned."),
    (BLUE, "NO PREDICTIVE DISPATCH",
     "Priority is decided by whoever takes the call and how panicked they sound. There is no data-driven assessment of how many officers are needed or how long an incident will last."),
]
for i, (col, title, body) in enumerate(problems):
    x = Inches(0.4) + i * Inches(4.28)
    add_rect(sl, x, Inches(1.5), Inches(4.0), Inches(4.5), fill=PANEL, line=col)
    add_rect(sl, x, Inches(1.5), Inches(4.0), Inches(0.06), fill=col)
    add_txt(sl, title, x + Inches(0.2), Inches(1.7), Inches(3.6), Inches(0.4),
            size=Pt(11), bold=True, color=col)
    add_txt(sl, body, x + Inches(0.2), Inches(2.2), Inches(3.6), Inches(3.6),
            size=Pt(11.5), color=GRAY, wrap=True)

# Bottom stat bar
add_rect(sl, 0, Inches(6.2), W, Inches(1.3), fill=RGBColor(0x0E, 0x10, 0x16))
stats2 = [
    ("10.2M", "Registered vehicles"),
    ("1,500+", "Added every day"),
    ("40+ min", "Avg Silk Board delay"),
    ("0", "Unified digital ops tools"),
]
for i, (n, l) in enumerate(stats2):
    x = Inches(1.2) + i * Inches(3.1)
    add_txt(sl, n, x, Inches(6.35), Inches(2), Inches(0.5),
            size=Pt(26), bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    add_txt(sl, l, x, Inches(6.85), Inches(2.5), Inches(0.3),
            size=Pt(9.5), color=GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SOLUTION OVERVIEW
# ════════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
amber_bar(sl)

section_chip(sl, "THE SOLUTION", Inches(0.4), Inches(0.2))
add_txt(sl, "DRISHTI — Three Layers, One Unified Platform",
        Inches(0.4), Inches(0.6), Inches(12), Inches(0.7),
        size=Pt(28), bold=True, color=WHITE)
add_rect(sl, Inches(0.4), Inches(1.25), Inches(4), Inches(0.04), fill=AMBER)

layers = [
    (CYAN,  "CITIZEN LAYER",
     "Public complaint portal — no login required. Citizen describes incident, gets DRS-BTP-XXXX tracking ID instantly. Real-time status timeline. AI-stamped priority score visible to reporter.",
     ["Multilingual (English + Kannada)", "GPS auto-detect or manual location", "Live status: Submitted → Deployed → Closed", "Firewall-rejected complaints flagged in red"]),
    (AMBER, "COMMAND CENTRE",
     "Admin and operator dashboards with live complaint triage, personnel map (Mappls SDK), deployment orders, and AI-generated dispatch recommendations for every incoming report.",
     ["AI priority score 0–100 per complaint", "Live Mappls map with officer GPS dots", "Approve/reject field access requests", "Real-time WebSocket ops chat"]),
    (GREEN, "FIELD OFFICER",
     "Mobile-responsive dashboard for officers on the ground. Shows active assignments, turn-by-turn routing (Mappls → OSRM fallback), GPS auto-polling every 30 seconds, and incident filing.",
     ["Auto GPS polling on login (30s cycle)", "Route navigation with API fallback chain", "File own incident reports from the field", "Real-time chat with Command Centre"]),
]
for i, (col, title, body, bullets) in enumerate(layers):
    x = Inches(0.3) + i * Inches(4.35)
    add_rect(sl, x, Inches(1.5), Inches(4.15), Inches(5.7), fill=PANEL, line=col)
    add_rect(sl, x, Inches(1.5), Inches(4.15), Inches(0.05), fill=col)
    add_txt(sl, title, x + Inches(0.15), Inches(1.65), Inches(3.8), Inches(0.3),
            size=Pt(10), bold=True, color=col)
    add_txt(sl, body, x + Inches(0.15), Inches(2.0), Inches(3.85), Inches(1.5),
            size=Pt(10.5), color=GRAY, wrap=True)
    for j, b in enumerate(bullets):
        add_txt(sl, f"→  {b}", x + Inches(0.15), Inches(3.6) + j * Inches(0.45),
                Inches(3.85), Inches(0.42), size=Pt(10.5), color=WHITE)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SYSTEM ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
amber_bar(sl)

section_chip(sl, "ARCHITECTURE", Inches(0.4), Inches(0.2))
add_txt(sl, "Full-Stack Architecture — End to End",
        Inches(0.4), Inches(0.6), Inches(10), Inches(0.6),
        size=Pt(26), bold=True, color=WHITE)
add_rect(sl, Inches(0.4), Inches(1.2), Inches(3.5), Inches(0.04), fill=AMBER)

# Row 1: Frontend
add_rect(sl, Inches(0.3), Inches(1.4), Inches(12.7), Inches(0.25), fill=RGBColor(0x0E, 0x10, 0x16))
add_txt(sl, "FRONTEND  —  Next.js 15 App Router  (Render Static)", Inches(0.5), Inches(1.4),
        Inches(12), Inches(0.25), size=Pt(9), bold=True, color=AMBER)

fe_boxes = ["Citizen Portal", "Admin Dashboard", "Operator Dashboard", "Field Officer Dashboard"]
for i, b in enumerate(fe_boxes):
    add_rect(sl, Inches(0.3) + i * Inches(3.2), Inches(1.7), Inches(3.0), Inches(0.6),
             fill=RGBColor(0x1A, 0x1F, 0x2E), line=CYAN)
    add_txt(sl, b, Inches(0.45) + i * Inches(3.2), Inches(1.82), Inches(2.8), Inches(0.4),
            size=Pt(10.5), bold=True, color=CYAN, align=PP_ALIGN.CENTER)

# Row 2: Backend
add_rect(sl, Inches(0.3), Inches(2.45), Inches(12.7), Inches(0.25), fill=RGBColor(0x0E, 0x10, 0x16))
add_txt(sl, "BACKEND  —  FastAPI + Uvicorn  (Render Web Service)", Inches(0.5), Inches(2.45),
        Inches(12), Inches(0.25), size=Pt(9), bold=True, color=AMBER)

be_boxes = [
    (AMBER, "Auth Service\nJWT + bcrypt\nHttpOnly Cookie"),
    (CYAN,  "Grievance API\nTriage Pipeline\nOutbox Events"),
    (GREEN, "Deployment API\nPersonnel GPS\nRoute Planning"),
    (BLUE,  "WebSocket Chat\nRoom Broadcast\nJWT ?token= auth"),
    (RED,   "ML Engine\nXGBoost / RF\nGemini Firewall"),
]
for i, (col, txt) in enumerate(be_boxes):
    add_rect(sl, Inches(0.3) + i * Inches(2.56), Inches(2.75), Inches(2.4), Inches(1.1),
             fill=RGBColor(0x10, 0x14, 0x1C), line=col)
    add_txt(sl, txt, Inches(0.4) + i * Inches(2.56), Inches(2.82), Inches(2.2), Inches(1.0),
            size=Pt(9.5), color=WHITE, wrap=True, align=PP_ALIGN.CENTER)

# Row 3: External services
add_rect(sl, Inches(0.3), Inches(4.0), Inches(12.7), Inches(0.25), fill=RGBColor(0x0E, 0x10, 0x16))
add_txt(sl, "EXTERNAL SERVICES  &  DATA LAYER", Inches(0.5), Inches(4.0),
        Inches(12), Inches(0.25), size=Pt(9), bold=True, color=AMBER)

ext = [
    (RGBColor(0x18,0x65,0xC8), "PostgreSQL\nSupabase\nFree Tier"),
    (RGBColor(0xC6,0x29,0x2B), "Redis\nUpstash TLS\nCache + Blacklist"),
    (RGBColor(0x1A,0x73,0xE8), "Gemini 2.0\nFlash Lite\nFirewall + Dispatch"),
    (RGBColor(0xFF,0x6F,0x00), "HuggingFace\nInference API\n384-dim Embed"),
    (RGBColor(0x00,0x96,0x88), "Mappls SDK\nGeocoding\n+ Routing"),
    (RGBColor(0x8B,0x5C,0xF6), "Resend API\nTransactional\nEmail"),
]
for i, (col, txt) in enumerate(ext):
    add_rect(sl, Inches(0.3) + i * Inches(2.14), Inches(4.3), Inches(2.0), Inches(0.9),
             fill=RGBColor(0x10, 0x14, 0x1C), line=col)
    add_txt(sl, txt, Inches(0.35) + i * Inches(2.14), Inches(4.35), Inches(1.9), Inches(0.85),
            size=Pt(9), color=WHITE, wrap=True, align=PP_ALIGN.CENTER)

# Key stats bottom
stats3 = [("FastAPI", "Backend"), ("Next.js 15", "Frontend"), ("psycopg3", "DB Driver"),
          ("PyJWT + bcrypt", "Auth"), ("Render", "Hosting")]
for i, (k, v) in enumerate(stats3):
    x = Inches(0.5) + i * Inches(2.56)
    add_txt(sl, k, x, Inches(5.45), Inches(2.3), Inches(0.3),
            size=Pt(11), bold=True, color=CYAN)
    add_txt(sl, v, x, Inches(5.75), Inches(2.3), Inches(0.3),
            size=Pt(9), color=GRAY)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — AI FIREWALL
# ════════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
amber_bar(sl)

section_chip(sl, "AI FIREWALL", Inches(0.4), Inches(0.2))
add_txt(sl, "Two-Gate Validation — Zero False Deletions",
        Inches(0.4), Inches(0.6), Inches(10), Inches(0.6),
        size=Pt(26), bold=True, color=WHITE)
add_rect(sl, Inches(0.4), Inches(1.2), Inches(3.8), Inches(0.04), fill=AMBER)

# Gate 1
add_rect(sl, Inches(0.3), Inches(1.45), Inches(6.0), Inches(2.8), fill=PANEL, line=CYAN)
add_rect(sl, Inches(0.3), Inches(1.45), Inches(6.0), Inches(0.06), fill=CYAN)
add_txt(sl, "GATE 1  ·  Keyword Pre-Filter", Inches(0.5), Inches(1.6),
        Inches(5.6), Inches(0.35), size=Pt(13), bold=True, color=CYAN)
add_txt(sl,
    "Zero latency. No network call. A frozen set of ~50 traffic keywords in English + transliterated Kannada scanned against the complaint description.\n\n"
    "English:  traffic, accident, junction, breakdown, collision, pothole, barricade, ambulance, flood, signal, lorry, convoy...\n\n"
    "Kannada:  jama, apagata, vahana, mara, raste, rasthe\n\n"
    "→  ANY match = PASS to Gate 2\n"
    "→  NO match = Instant reject (not a traffic incident)",
    Inches(0.5), Inches(2.05), Inches(5.6), Inches(2.0),
    size=Pt(10.5), color=GRAY, wrap=True)

# Arrow
add_txt(sl, "PASS  ▶", Inches(6.45), Inches(2.65), Inches(1.2), Inches(0.4),
        size=Pt(13), bold=True, color=AMBER, align=PP_ALIGN.CENTER)

# Gate 2
add_rect(sl, Inches(7.3), Inches(1.45), Inches(5.7), Inches(2.8), fill=PANEL, line=AMBER)
add_rect(sl, Inches(7.3), Inches(1.45), Inches(5.7), Inches(0.06), fill=AMBER)
add_txt(sl, "GATE 2  ·  Gemini 2.0 Flash Lite", Inches(7.5), Inches(1.6),
        Inches(5.3), Inches(0.35), size=Pt(13), bold=True, color=AMBER)
add_txt(sl,
    "Runs ASYNCHRONOUSLY after HTTP response sent — citizen already has their tracking ID.\n\n"
    'Prompt:  "Is this a road traffic incident? Reply YES or NO."\n\n'
    "→  Explicit NO  =  Flag [FIREWALL REJECTED] on record\n"
    "→  YES  =  Complaint stays clean\n"
    "→  Timeout / Error / Ambiguous  =  PASS THROUGH\n\n"
    "20-second timeout. Only explicit rejection causes action.\nCitizen is never penalised for Gemini being slow.",
    Inches(7.5), Inches(2.05), Inches(5.3), Inches(2.0),
    size=Pt(10.5), color=GRAY, wrap=True)

# Why this design
add_rect(sl, Inches(0.3), Inches(4.4), Inches(12.7), Inches(2.8), fill=RGBColor(0x0C, 0x10, 0x18), line=BORDER)
add_txt(sl, "WHY THIS DESIGN", Inches(0.5), Inches(4.55),
        Inches(5), Inches(0.3), size=Pt(11), bold=True, color=AMBER)

reasons = [
    (GREEN,  "Fail-Safe",       "Errors and timeouts PASS complaints through. We never wrongly delete a legitimate truck collision at 11 PM because Gemini had a slow second."),
    (CYAN,   "Zero Latency UX", "Gate 1 is pure Python — sub-millisecond. The citizen response time is never held hostage by an external API call."),
    (AMBER,  "Async Firewall",  "Gemini runs in a FastAPI BackgroundTask. The citizen already has their DRS-BTP-XXXX ID before Gemini even starts thinking."),
    (BLUE,   "Multilingual",    "Keyword set includes transliterated Kannada because a significant fraction of Bengaluru reporters describe incidents in Kannada, not English."),
]
for i, (col, title, body) in enumerate(reasons):
    x = Inches(0.5) + i * Inches(3.18)
    add_txt(sl, f"▸  {title}", x, Inches(4.9), Inches(3.0), Inches(0.3),
            size=Pt(10), bold=True, color=col)
    add_txt(sl, body, x, Inches(5.25), Inches(3.0), Inches(1.7),
            size=Pt(9.5), color=GRAY, wrap=True)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ML PIPELINE
# ════════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
amber_bar(sl)

section_chip(sl, "ML PIPELINE", Inches(0.4), Inches(0.2))
add_txt(sl, "Three-Level Cascade — Always Returns a Result",
        Inches(0.4), Inches(0.6), Inches(10), Inches(0.6),
        size=Pt(26), bold=True, color=WHITE)
add_rect(sl, Inches(0.4), Inches(1.2), Inches(4.2), Inches(0.04), fill=AMBER)

# Level boxes
levels = [
    (CYAN,  "LEVEL 1  —  XGBoost (Primary)",
     [("duration_model.pkl",  "Predicts incident resolution time in minutes"),
      ("resource_model.pkl",  "Binary priority — High (1) or Low (0)"),
      ("397 features",        "13 structural + 384 HuggingFace embedding dims"),
      ("Personnel calc",      "Rule-derived from duration + priority + road closure"),
      ("Urgency",             "CRITICAL / HIGH / MEDIUM / LOW from combined output")]),
    (AMBER, "LEVEL 2  —  Random Forest (Fallback)",
     [("impact_model.pkl",    "Classifies severity: Critical / High / Medium / Low"),
      ("traffic_duration_rf", "Regresses incident duration in minutes"),
      ("Simpler features",    "No embeddings — structured signals only"),
      ("Auto-activates",      "When XGBoost fails to load or predicts incorrectly"),
      ("Same output schema",  "Drops into same recommendation builder seamlessly")]),
    (GREEN, "LEVEL 3  —  Rule-Based (Last Resort)",
     [("Keyword severity",    "Ambulance/death → Critical  |  Blocked → High"),
      ("Type defaults",       "Accident/Road closure → High  |  Parking → Low"),
      ("Score lookup",        "Severity → fixed priority score 0–100"),
      ("4 text templates",    "90+ / 75–90 / 50–75 / 0–50 score band responses"),
      ("Zero dependencies",   "Runs with no model files, no API keys, no RAM")]),
]
for i, (col, title, items) in enumerate(levels):
    x = Inches(0.3) + i * Inches(4.35)
    add_rect(sl, x, Inches(1.45), Inches(4.15), Inches(4.3), fill=PANEL, line=col)
    add_rect(sl, x, Inches(1.45), Inches(4.15), Inches(0.05), fill=col)
    add_txt(sl, title, x + Inches(0.15), Inches(1.6), Inches(3.85), Inches(0.32),
            size=Pt(10.5), bold=True, color=col)
    for j, (k, v) in enumerate(items):
        y = Inches(2.05) + j * Inches(0.68)
        add_txt(sl, k, x + Inches(0.15), y, Inches(1.65), Inches(0.3),
                size=Pt(9.5), bold=True, color=WHITE)
        add_txt(sl, v, x + Inches(1.85), y, Inches(2.2), Inches(0.55),
                size=Pt(9.5), color=GRAY, wrap=True)

# Bottom — Gemini dispatch
add_rect(sl, Inches(0.3), Inches(5.9), Inches(12.7), Inches(1.4), fill=RGBColor(0x0C, 0x10, 0x18), line=BORDER)
add_txt(sl, "AFTER ML  →  Gemini 2.0 Flash Lite generates natural language dispatch recommendation",
        Inches(0.5), Inches(6.0), Inches(12.2), Inches(0.35),
        size=Pt(11), bold=True, color=AMBER)
add_txt(sl,
    '"Dispatch 1 Inspector, 2 Sub-Inspectors and 6 Constables to Silk Board Junction (ORR East 2 corridor), South Zone — priority score 84/100, truck breakdown blocking 2 lanes."  '
    ' |  Sentence 2 covers ETA, BBMP coordination, tow unit, ambulance if needed.',
    Inches(0.5), Inches(6.4), Inches(12.2), Inches(0.75),
    size=Pt(10), color=GRAY, wrap=True, italic=True)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — WHY THESE MODELS
# ════════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
amber_bar(sl)

section_chip(sl, "MODEL JUSTIFICATION", Inches(0.4), Inches(0.2))
add_txt(sl, "Why XGBoost + MiniLM — Not Deep Learning End-to-End",
        Inches(0.4), Inches(0.6), Inches(12), Inches(0.6),
        size=Pt(26), bold=True, color=WHITE)
add_rect(sl, Inches(0.4), Inches(1.2), Inches(5.5), Inches(0.04), fill=AMBER)

cards = [
    (CYAN,  "XGBoost for Tabular Data",
     "Traffic incidents are fundamentally tabular problems — time of day, GPS coordinates, corridor, zone, vehicle type. XGBoost is the gold standard for structured data with mixed categorical and numerical features. It handles missing values natively, is interpretable (feature importances), trains fast, and serializes to a 2MB file. A neural network would need 10x the data and 100x the RAM for no gain on this kind of structured signal.",
     ["Handles mixed types natively", "Serializes to ~2MB (fits on free tier)", "Feature importance = explainable AI", "Proven SOTA on tabular benchmarks"]),
    (AMBER, "MiniLM for Text Embeddings",
     "paraphrase-multilingual-MiniLM-L12-v2 is a distilled sentence transformer — 22M parameters vs. 110M for BERT. It handles English AND Kannada (plus 50 other languages) because Bengaluru reporters write in both. The 384-dim output is compact enough that XGBoost can consume it as features without a separate fusion layer. Critical choice: we call it via HuggingFace Inference API — zero RAM on our server.",
     ["Multilingual: English + Kannada", "384 dims → direct XGBoost features", "Zero server RAM (HF Inference API)", "Same model used for training and inference"]),
    (GREEN, "Gemini 2.0 Flash Lite for NLG",
     "Dispatch recommendations must be in operational police language, location-specific, and rank-aware. A template system produces generic output. Gemini Flash Lite is the most cost-effective Gemini model — minimal latency, high quality, free tier quota sufficient for a hackathon volume. It's constrained to exactly 2 sentences with a strict prompt so it can't hallucinate long paragraphs.",
     ["2-sentence structured output", "Location + rank + duration aware", "BTP operational language trained in", "Falls back to templates if unavailable"]),
]
for i, (col, title, body, tags) in enumerate(cards):
    x = Inches(0.3) + i * Inches(4.35)
    add_rect(sl, x, Inches(1.5), Inches(4.15), Inches(5.7), fill=PANEL, line=col)
    add_rect(sl, x, Inches(1.5), Inches(4.15), Inches(0.05), fill=col)
    add_txt(sl, title, x + Inches(0.15), Inches(1.65),
            Inches(3.85), Inches(0.32), size=Pt(11), bold=True, color=col)
    add_txt(sl, body, x + Inches(0.15), Inches(2.05),
            Inches(3.85), Inches(2.5), size=Pt(10), color=GRAY, wrap=True)
    for j, tag in enumerate(tags):
        ty = Inches(4.7) + j * Inches(0.38)
        add_rect(sl, x + Inches(0.15), ty, Inches(3.75), Inches(0.32),
                 fill=RGBColor(0x10, 0x14, 0x1C), line=col)
        add_txt(sl, f"✓  {tag}", x + Inches(0.28), ty + Inches(0.04),
                Inches(3.5), Inches(0.28), size=Pt(9.5), color=col)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — FEATURE ENGINEERING
# ════════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
amber_bar(sl)

section_chip(sl, "FEATURE ENGINEERING", Inches(0.4), Inches(0.2))
add_txt(sl, "397 Features — What the Model Actually Sees",
        Inches(0.4), Inches(0.6), Inches(10), Inches(0.6),
        size=Pt(26), bold=True, color=WHITE)
add_rect(sl, Inches(0.4), Inches(1.2), Inches(3.8), Inches(0.04), fill=AMBER)

# Structural features
add_rect(sl, Inches(0.3), Inches(1.45), Inches(5.6), Inches(5.8), fill=PANEL, line=CYAN)
add_rect(sl, Inches(0.3), Inches(1.45), Inches(5.6), Inches(0.05), fill=CYAN)
add_txt(sl, "STRUCTURAL FEATURES  (13)", Inches(0.45), Inches(1.6),
        Inches(5.3), Inches(0.3), size=Pt(11), bold=True, color=CYAN)

struct = [
    ("hour", "0–23 — captures morning/evening rush"),
    ("day_of_week", "0–6 — weekday vs weekend patterns"),
    ("month", "1–12 — seasonal flooding, events"),
    ("is_night", "1 if hour ≥ 21 or ≤ 6 — severity proxy"),
    ("reporting_delay_min", "Time from incident to report"),
    ("latitude / longitude", "GPS — corridor and zone proxy"),
    ("requires_road_closure", "Binary — keyword + type derived"),
    ("event_cause_enc", "9 classes: accident, tree fall, flood…"),
    ("veh_type_enc", "6 classes: heavy, bus, two-wheeler…"),
    ("corridor_enc", "20+ Bengaluru corridors: ORR, MG Rd…"),
    ("police_station_enc", "Jurisdiction encoding"),
    ("zone_enc", "Central / N / S / E / W + subdivisions"),
]
for i, (feat, desc) in enumerate(struct):
    y = Inches(2.0) + i * Inches(0.38)
    add_txt(sl, feat, Inches(0.45), y, Inches(1.9), Inches(0.35),
            size=Pt(9.5), bold=True, color=WHITE)
    add_txt(sl, desc, Inches(2.38), y, Inches(3.3), Inches(0.35),
            size=Pt(9.5), color=GRAY, wrap=True)

# Embedding features
add_rect(sl, Inches(6.1), Inches(1.45), Inches(6.9), Inches(5.8), fill=PANEL, line=AMBER)
add_rect(sl, Inches(6.1), Inches(1.45), Inches(6.9), Inches(0.05), fill=AMBER)
add_txt(sl, "TEXT EMBEDDING FEATURES  (384)", Inches(6.25), Inches(1.6),
        Inches(6.6), Inches(0.3), size=Pt(11), bold=True, color=AMBER)

add_txt(sl, "Model:  sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2",
        Inches(6.25), Inches(2.0), Inches(6.6), Inches(0.3),
        size=Pt(10), bold=True, color=WHITE)
add_txt(sl,
    "Each complaint description is encoded into a 384-dimensional dense vector. "
    "Every dimension captures a different semantic axis of the text — from \"is this urgent?\" "
    "to \"does this involve a heavy vehicle?\" to \"is flooding mentioned?\"\n\n"
    "The MiniLM architecture is a distilled version of BERT (22M vs 110M parameters). "
    "It was trained on 1 billion sentence pairs across 50 languages including Kannada, "
    "making it ideal for Bengaluru's bilingual reporting pattern.\n\n"
    "We call it via the HuggingFace Inference API — the model never loads into our "
    "server's memory. On Render's 512MB free tier, this is non-negotiable. The same "
    "model was used during training AND inference, so there is zero embedding drift.",
    Inches(6.25), Inches(2.4), Inches(6.6), Inches(2.5),
    size=Pt(10.5), color=GRAY, wrap=True)

add_txt(sl, "emb_0, emb_1, emb_2  ...  emb_383",
        Inches(6.25), Inches(4.95), Inches(6.6), Inches(0.35),
        size=Pt(11), bold=True, color=AMBER)

# NLP extraction note
add_rect(sl, Inches(6.25), Inches(5.35), Inches(6.6), Inches(1.7), fill=RGBColor(0x10, 0x14, 0x1C), line=BORDER)
add_txt(sl, "NLP PRE-EXTRACTION  (before encoding)",
        Inches(6.4), Inches(5.45), Inches(6.3), Inches(0.3),
        size=Pt(9.5), bold=True, color=CYAN)
add_txt(sl,
    "Event cause extracted by keyword scan (9 categories).\n"
    "Vehicle type extracted by keyword scan (6 categories).\n"
    "These feed the structural encoders as well as the embedding input.",
    Inches(6.4), Inches(5.8), Inches(6.3), Inches(1.1),
    size=Pt(9.5), color=GRAY, wrap=True)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — SCALABILITY
# ════════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
amber_bar(sl)

section_chip(sl, "SCALABILITY", Inches(0.4), Inches(0.2))
add_txt(sl, "How DRISHTI Scales — From Hackathon to City-Wide",
        Inches(0.4), Inches(0.6), Inches(12), Inches(0.6),
        size=Pt(26), bold=True, color=WHITE)
add_rect(sl, Inches(0.4), Inches(1.2), Inches(4.5), Inches(0.04), fill=AMBER)

# Phase cards
phases = [
    (CYAN, "PHASE 1  ·  NOW\nHackathon / Pilot",
     ["Single Render free-tier backend", "Single PostgreSQL (Supabase free)", "Upstash Redis (free tier)", "HuggingFace Inference API (shared)", "~50 concurrent users", "~500 grievances/day capacity"],
     "Render Free · Supabase Free · Upstash Free"),
    (AMBER, "PHASE 2  ·  DISTRICT\nBTP Pilot Deployment",
     ["Render paid tier (4GB RAM, autoscale)", "Supabase Pro (connection pooling)", "Redis cluster (dedicated)", "Multiple Uvicorn workers (gunicorn)", "Batch HuggingFace or local MiniLM", "~500 concurrent · ~10K grievances/day"],
     "~₹8K/month infrastructure"),
    (GREEN, "PHASE 3  ·  CITY-WIDE\nAll BTP Zones",
     ["Kubernetes cluster (GCP/AWS India)", "PostgreSQL read replicas per zone", "Redis Cluster with horizontal sharding", "Model serving via TorchServe/TF Serving", "Kafka for async event streaming", "100K+ grievances/day · 5K concurrent"],
     "Enterprise cloud · ₹80K-2L/month"),
]
for i, (col, title, bullets, cost) in enumerate(phases):
    x = Inches(0.3) + i * Inches(4.35)
    add_rect(sl, x, Inches(1.5), Inches(4.15), Inches(4.8), fill=PANEL, line=col)
    add_rect(sl, x, Inches(1.5), Inches(4.15), Inches(0.05), fill=col)
    add_txt(sl, title, x + Inches(0.15), Inches(1.65), Inches(3.85), Inches(0.55),
            size=Pt(10.5), bold=True, color=col)
    for j, b in enumerate(bullets):
        add_txt(sl, f"·  {b}", x + Inches(0.15), Inches(2.3) + j * Inches(0.42),
                Inches(3.85), Inches(0.4), size=Pt(10), color=GRAY)
    add_rect(sl, x + Inches(0.15), Inches(5.8), Inches(3.85), Inches(0.3),
             fill=RGBColor(0x10, 0x14, 0x1C), line=col)
    add_txt(sl, cost, x + Inches(0.25), Inches(5.82), Inches(3.65), Inches(0.28),
            size=Pt(9), bold=True, color=col)

# Architecture advantages for scale
add_rect(sl, Inches(0.3), Inches(6.45), Inches(12.7), Inches(0.85), fill=RGBColor(0x0C, 0x10, 0x18), line=BORDER)
add_txt(sl, "BUILT-IN SCALE ADVANTAGES:", Inches(0.5), Inches(6.52),
        Inches(3), Inches(0.3), size=Pt(9.5), bold=True, color=AMBER)
scalepts = [
    "Transactional outbox → no event loss on crash",
    "Redis cache → DB queries drop 80%+ at scale",
    "Stateless FastAPI workers → horizontal scaling trivial",
    "HuggingFace API → swap to local model with one env var",
    "Three-level ML fallback → system never fails silently",
]
for i, pt in enumerate(scalepts):
    add_txt(sl, f"✓  {pt}", Inches(3.6) + i * Inches(1.92), Inches(6.52),
            Inches(1.88), Inches(0.65), size=Pt(8.5), color=GRAY, wrap=True)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — IMPACT & VISION
# ════════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
amber_bar(sl)

section_chip(sl, "IMPACT & VISION", Inches(0.4), Inches(0.2))
add_txt(sl, "Real Impact — Measurable Outcomes",
        Inches(0.4), Inches(0.6), Inches(10), Inches(0.6),
        size=Pt(26), bold=True, color=WHITE)
add_rect(sl, Inches(0.4), Inches(1.2), Inches(3.5), Inches(0.04), fill=AMBER)

# Impact metrics
impacts = [
    (CYAN,  "< 2 sec",  "Complaint intake to tracking ID issued to citizen"),
    (GREEN, "397 dims", "Feature vector per complaint — richer than any manual triage"),
    (AMBER, "3-level",  "ML cascade — system always returns a result, never blank"),
    (RED,   "0 RAM",    "For HF embeddings — API call, not local model load"),
    (BLUE,  "30 sec",   "GPS refresh cycle — live officer positions on command map"),
    (CYAN,  "20 sec",   "Gemini firewall timeout before auto-pass — no user delay"),
]
for i, (col, num, label) in enumerate(impacts):
    row, col_i = divmod(i, 3)
    x = Inches(0.3) + col_i * Inches(4.35)
    y = Inches(1.5) + row * Inches(1.1)
    add_rect(sl, x, y, Inches(4.15), Inches(0.95), fill=PANEL, line=col)
    add_txt(sl, num, x + Inches(0.2), y + Inches(0.08),
            Inches(1.6), Inches(0.45), size=Pt(22), bold=True, color=col)
    add_txt(sl, label, x + Inches(1.85), y + Inches(0.15),
            Inches(2.2), Inches(0.65), size=Pt(9.5), color=GRAY, wrap=True)

# What this replaces
add_rect(sl, Inches(0.3), Inches(3.8), Inches(6.2), Inches(3.4), fill=PANEL, line=RED)
add_rect(sl, Inches(0.3), Inches(3.8), Inches(6.2), Inches(0.05), fill=RED)
add_txt(sl, "WHAT DRISHTI REPLACES", Inches(0.5), Inches(3.95),
        Inches(5.8), Inches(0.3), size=Pt(11), bold=True, color=RED)
befores = [
    "Phone calls to control room with no log or timestamp",
    "WhatsApp forwarding chains for incident reports",
    "Radio dispatch with no digital record",
    "Manual GPS tracking via phone call check-ins",
    "Gut-feel priority assignment by duty officer",
    "No citizen feedback loop after filing complaint",
]
for i, b in enumerate(befores):
    add_txt(sl, f"✗  {b}", Inches(0.5), Inches(4.35) + i * Inches(0.38),
            Inches(5.8), Inches(0.35), size=Pt(10), color=RGBColor(0xFC, 0xA5, 0xA5))

# Future roadmap
add_rect(sl, Inches(6.7), Inches(3.8), Inches(6.3), Inches(3.4), fill=PANEL, line=GREEN)
add_rect(sl, Inches(6.7), Inches(3.8), Inches(6.3), Inches(0.05), fill=GREEN)
add_txt(sl, "NEXT — ROADMAP", Inches(6.9), Inches(3.95),
        Inches(5.8), Inches(0.3), size=Pt(11), bold=True, color=GREEN)
nexts = [
    ("CCTV Integration", "Feed live camera frames to vision model for auto-detection"),
    ("Predictive Hotspots", "Retrain on 90-day rolling data — predict junctions at risk"),
    ("WhatsApp Bot", "Twilio integration — citizens report via existing WhatsApp"),
    ("Kannada Voice Input", "STT → transliterate → same ML pipeline"),
    ("Continuous Retraining", "Feedback loop: closed incidents label new training rows"),
    ("Multi-City", "Mysuru, Hubballi, Chennai — same platform, city-specific models"),
]
for i, (title, desc) in enumerate(nexts):
    y = Inches(4.35) + i * Inches(0.38)
    add_txt(sl, f"▸  {title}", Inches(6.9), y, Inches(2.1), Inches(0.35),
            size=Pt(9.5), bold=True, color=GREEN)
    add_txt(sl, desc, Inches(9.1), y, Inches(3.8), Inches(0.35),
            size=Pt(9.5), color=GRAY, wrap=True)

# Final tagline
add_rect(sl, 0, Inches(6.8), W, Inches(0.7), fill=RGBColor(0x0E, 0x10, 0x16))
add_txt(sl,
    "DRISHTI doesn't replace the officer's judgement — it gives that officer better information, faster.  "
    "In a city of 10 million vehicles, shaving 15 minutes off average response time saves lives.",
    Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.55),
    size=Pt(11), color=AMBER, align=PP_ALIGN.CENTER, italic=True)

# ── Save ────────────────────────────────────────────────────────────────────────
out = r"c:\Users\bhard\Desktop\DRISHTI_Presentation.pptx"
prs.save(out)
print(f"Saved to: {out}")
